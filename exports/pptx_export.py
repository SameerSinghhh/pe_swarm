"""
Professional PowerPoint deck generator.

Takes pre-computed analysis + research + value creation data and renders
it into a polished, board-ready presentation. No computation here —
purely rendering pre-verified numbers into formatted slides.
"""

from datetime import date
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

from analysis.types import AnalysisResult, Favorability
from research.types import ResearchBrief
from value_creation.types import ValueCreationPlan
from modeling.types import ModelResult


# ── Colors ──
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x00, 0x61, 0x00)
RED = RGBColor(0x9C, 0x00, 0x06)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


# ── Format helpers ──

def _fmt_dollar(val, short=True):
    if val is None: return "—"
    if short:
        if abs(val) >= 1e9: return f"${val/1e9:.1f}B"
        if abs(val) >= 1e6: return f"${val/1e6:.1f}M"
        if abs(val) >= 1e3: return f"${val/1e3:.0f}K"
    return f"${val:,.0f}"

def _fmt_dollar_signed(val):
    if val is None: return "—"
    if val >= 0: return f"+{_fmt_dollar(val)}"
    return f"-{_fmt_dollar(abs(val))}"

def _fmt_pct(val):
    if val is None: return "—"
    return f"{val:.1f}%"

def _fmt_ratio(val):
    if val is None: return "—"
    return f"{val:.1f}x"


# ── Slide helpers ──

def _set_cell(cell, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, fill_color=None):
    """Set cell text with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def _add_table(slide, data, left, top, width, col_widths, font_size=11):
    """Add a professional table to a slide. First row is header."""
    rows = len(data)
    cols = len(data[0]) if data else 0
    if rows == 0 or cols == 0: return None

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.35 * rows))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths[:cols]):
        table.columns[i].width = Inches(w)

    # Header row
    for j, val in enumerate(data[0]):
        _set_cell(table.cell(0, j), val, font_size=font_size, bold=True,
                  color=WHITE, alignment=PP_ALIGN.CENTER, fill_color=DARK_BLUE)

    # Data rows
    for i in range(1, rows):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        for j, val in enumerate(data[i]):
            color = BLACK
            # Color positive/negative values
            if isinstance(val, str) and val.startswith("+$"):
                color = GREEN
            elif isinstance(val, str) and val.startswith("-$"):
                color = RED
            elif isinstance(val, str) and val.startswith("($"):
                color = RED
            _set_cell(table.cell(i, j), val, font_size=font_size,
                      color=color, fill_color=bg)

    return table


def _add_title(slide, title, subtitle=""):
    """Add title and optional subtitle to slide."""
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    run.font.name = "Calibri"

    if subtitle:
        top2 = Inches(0.85)
        txBox2 = slide.shapes.add_textbox(left, top2, Inches(9), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = GRAY
        run2.font.name = "Calibri"


def _add_metric_row(slide, metrics, top):
    """Add a row of key metrics (label + value pairs)."""
    n = len(metrics)
    if n == 0: return
    col_width = 9.0 / n

    for i, (label, value) in enumerate(metrics):
        left = Inches(0.5 + i * col_width)
        # Value
        txBox = slide.shapes.add_textbox(left, top, Inches(col_width - 0.1), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(value)
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = DARK_BLUE
        run.font.name = "Calibri"
        # Label
        txBox2 = slide.shapes.add_textbox(left, Inches(top.inches + 0.5), Inches(col_width - 0.1), Inches(0.3))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = label
        run2.font.size = Pt(10)
        run2.font.color.rgb = GRAY
        run2.font.name = "Calibri"


def _add_text_block(slide, text, left, top, width, height, font_size=12, color=BLACK):
    """Add a text block."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)[:1500]  # Truncate for safety
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_bullet_list(slide, items, left, top, width, height, font_size=12):
    """Add a bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:8]):  # Max 8 items per slide
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        # Clean the text
        clean = str(item).replace("{", "").replace("}", "")[:200]
        run.text = f"• {clean}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ═══════════════════════════════════════════════════════════════════════════

def export_to_pptx(
    analysis: AnalysisResult,
    research_brief: Optional[ResearchBrief] = None,
    value_creation: Optional[ValueCreationPlan] = None,
    model_result: Optional[ModelResult] = None,
    company_name: str = "",
    sector: str = "",
    filepath=None,
):
    """Generate a professional PowerPoint deck. All numbers are pre-computed."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    _create_cover(prs, company_name, sector)

    # Slide 2: Executive Summary
    _create_summary(prs, analysis, value_creation, model_result)

    # Slide 3: Historical Performance
    if analysis.margins:
        _create_performance(prs, analysis)

    # Slide 4: EBITDA Bridge
    if analysis.ebitda_bridges:
        _create_bridge(prs, analysis.ebitda_bridges)

    # Slide 5: Variance
    if analysis.variance:
        _create_variance(prs, analysis.variance)

    # Slide 6: Peer Benchmarking
    if research_brief and research_brief.peer_companies:
        _create_peers(prs, research_brief, analysis)

    # Slide 7: Gap Analysis
    if research_brief and research_brief.gaps:
        _create_gaps(prs, research_brief)

    # Slide 8: Working Capital + FCF
    if analysis.working_capital or analysis.fcf:
        _create_wc_fcf(prs, analysis)

    # Slide 9: Value Creation Plan
    if value_creation and value_creation.prioritized_plan:
        _create_vc_plan(prs, value_creation)

    # Slide 10: AI Roadmap
    if value_creation and (value_creation.ai_automation_opportunities or value_creation.ai_product_recommendations):
        _create_ai_roadmap(prs, value_creation)

    # Slide 11: Strategic + Risks
    if value_creation and (value_creation.strategic_priorities or value_creation.key_risks):
        _create_strategy_risks(prs, value_creation)

    # Slide 12: Returns
    if model_result and model_result.returns:
        _create_returns(prs, model_result.returns)

    prs.save(filepath)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def _create_cover(prs, company, sector):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Blue bar at top
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Company name
    _add_text_block(slide, company, Inches(1), Inches(2.2), Inches(11), Inches(1.2), font_size=44, color=DARK_BLUE)

    # Subtitle
    _add_text_block(slide, "Value Creation Analysis", Inches(1), Inches(3.4), Inches(11), Inches(0.6), font_size=24, color=GRAY)

    # Sector + date
    _add_text_block(slide, f"{sector}  •  {date.today().strftime('%B %Y')}", Inches(1), Inches(4.2), Inches(11), Inches(0.4), font_size=14, color=GRAY)

    # Bottom bar
    shape2 = slide.shapes.add_shape(1, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = DARK_BLUE
    shape2.line.fill.background()


def _create_summary(prs, analysis, vc_plan, model):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Executive Summary")

    metrics = []
    if analysis.ltm:
        l = analysis.ltm
        if l.ltm_revenue: metrics.append(("LTM Revenue", _fmt_dollar(l.ltm_revenue)))
        if l.ltm_ebitda: metrics.append(("LTM EBITDA", _fmt_dollar(l.ltm_ebitda)))
        if l.ltm_ebitda_margin_pct: metrics.append(("EBITDA Margin", _fmt_pct(l.ltm_ebitda_margin_pct)))
        if l.rule_of_40: metrics.append(("Rule of 40", f"{l.rule_of_40:.0f}"))
    if model and model.returns:
        if model.returns.moic: metrics.append(("MOIC", _fmt_ratio(model.returns.moic)))
        if model.returns.irr: metrics.append(("IRR", f"{model.returns.irr:.0%}"))

    if metrics:
        _add_metric_row(slide, metrics[:6], Inches(1.3))

    # Total opportunity
    if vc_plan and vc_plan.total_ebitda_opportunity > 0:
        _add_text_block(slide, f"Total EBITDA Opportunity: {_fmt_dollar(vc_plan.total_ebitda_opportunity)}/year",
                       Inches(0.5), Inches(2.5), Inches(12), Inches(0.5), font_size=20, color=GREEN)

    # Executive summary text
    if vc_plan and vc_plan.executive_summary:
        import re
        clean = re.sub(r'\$[^$]*\$', '', vc_plan.executive_summary)
        clean = re.sub(r'[{}]', '', clean)
        _add_text_block(slide, clean[:800], Inches(0.5), Inches(3.2), Inches(12), Inches(3.5), font_size=13, color=DARK_GRAY)


def _create_performance(prs, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Historical Performance", "Margin & growth trends")

    periods = analysis.margins.periods[-6:]  # Last 6
    header = ["Metric"] + [p.period for p in periods]
    rows = [header]

    for label, attr in [
        ("Gross Margin", "gross_margin_pct"),
        ("EBITDA Margin", "ebitda_margin_pct"),
        ("S&M % Revenue", "sm_pct_revenue"),
        ("R&D % Revenue", "rd_pct_revenue"),
        ("G&A % Revenue", "ga_pct_revenue"),
        ("Rev Growth MoM", "revenue_growth_mom"),
        ("Rev Growth YoY", "revenue_growth_yoy"),
    ]:
        row = [label]
        for p in periods:
            val = getattr(p, attr, None)
            row.append(_fmt_pct(val) if val is not None else "—")
        rows.append(row)

    col_widths = [1.8] + [1.5] * len(periods)
    _add_table(slide, rows, Inches(0.5), Inches(1.5), Inches(12), col_widths, font_size=10)


def _create_bridge(prs, bridges):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "EBITDA Bridge")

    b = bridges.mom
    if not b: return

    header = ["Component", "Amount"]
    rows = [header]
    rows.append(["Starting EBITDA (" + b.base_period + ")", _fmt_dollar(b.base_ebitda, short=False)])
    for c in b.components:
        rows.append(["  " + c.name, _fmt_dollar_signed(c.value)])
    rows.append(["Ending EBITDA (" + b.current_period + ")", _fmt_dollar(b.current_ebitda, short=False)])
    rows.append(["Total Change", _fmt_dollar_signed(b.total_change)])

    _add_table(slide, rows, Inches(0.5), Inches(1.5), Inches(5.5), [3.0, 2.5])

    # PY bridge next to it if available
    if bridges.vs_prior_year:
        py = bridges.vs_prior_year
        rows2 = [["Component", "Amount"]]
        rows2.append(["Starting (" + py.base_period + ")", _fmt_dollar(py.base_ebitda, short=False)])
        for c in py.components:
            rows2.append(["  " + c.name, _fmt_dollar_signed(c.value)])
        rows2.append(["Ending (" + py.current_period + ")", _fmt_dollar(py.current_ebitda, short=False)])
        rows2.append(["Total Change", _fmt_dollar_signed(py.total_change)])
        _add_table(slide, rows2, Inches(6.8), Inches(1.5), Inches(5.5), [3.0, 2.5])


def _create_variance(prs, variance):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    latest = variance.periods[-1]
    _add_title(slide, f"Variance Analysis — {latest.period}")

    data = latest.vs_prior_month
    if not data: return

    header = ["Line Item", "Actual", "Prior Month", "$ Change", "% Change", ""]
    rows = [header]
    for v in data:
        fav = "✓" if v.favorable == Favorability.FAVORABLE else ("✗" if v.favorable == Favorability.UNFAVORABLE else "—")
        rows.append([
            v.line_item.replace("_", " ").title(),
            _fmt_dollar(v.actual, short=False),
            _fmt_dollar(v.comparator, short=False),
            _fmt_dollar_signed(v.dollar_change),
            _fmt_pct(v.pct_change) if v.pct_change else "—",
            fav,
        ])

    _add_table(slide, rows, Inches(0.5), Inches(1.5), Inches(12), [2.2, 1.8, 1.8, 1.8, 1.2, 0.6])


def _create_peers(prs, research, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Peer Benchmarking", "Public comparable companies")

    header = ["Company", "Revenue", "Gross Margin", "EBITDA Margin", "Growth YoY", "EV/EBITDA"]
    rows = [header]

    # Our company first
    if analysis.ltm and analysis.margins and analysis.margins.periods:
        m = analysis.margins.periods[-1]
        rows.append([
            f"★ {research.company_name}",
            _fmt_dollar(analysis.ltm.ltm_revenue),
            _fmt_pct(m.gross_margin_pct),
            _fmt_pct(m.ebitda_margin_pct),
            _fmt_pct(m.revenue_growth_yoy),
            "—",
        ])

    for p in research.peer_companies[:5]:
        rows.append([
            f"{p.name} ({p.ticker})",
            _fmt_dollar(p.revenue),
            _fmt_pct(p.gross_margin_pct),
            _fmt_pct(p.ebitda_margin_pct),
            _fmt_pct(p.revenue_growth_yoy_pct),
            _fmt_ratio(p.ev_to_ebitda),
        ])

    _add_table(slide, rows, Inches(0.5), Inches(1.5), Inches(12), [3.5, 1.5, 1.5, 1.5, 1.5, 1.5])


def _create_gaps(prs, research):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Gap Analysis", "Company vs peer median")

    header = ["Metric", "Company", "Peer Median", "Gap", "Status"]
    rows = [header]
    for g in research.gaps:
        status = "Strength" if g.gap > 2 else ("Gap" if g.gap < -2 else "In Line")
        rows.append([
            g.metric,
            _fmt_pct(g.company_value),
            _fmt_pct(g.peer_median),
            f"{g.gap:+.1f}pp",
            status,
        ])

    _add_table(slide, rows, Inches(0.5), Inches(1.5), Inches(10), [2.5, 1.5, 1.5, 1.5, 1.5])


def _create_wc_fcf(prs, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Working Capital & Cash Flow")

    metrics = []
    if analysis.working_capital and analysis.working_capital.periods:
        wc = analysis.working_capital.periods[-1]
        if wc.dso: metrics.append(("DSO", f"{wc.dso:.0f} days"))
        if wc.dpo: metrics.append(("DPO", f"{wc.dpo:.0f} days"))
        if wc.ccc: metrics.append(("CCC", f"{wc.ccc:.0f} days"))

    if analysis.fcf and analysis.fcf.periods:
        f = analysis.fcf.periods[-1]
        if f.free_cash_flow: metrics.append(("FCF", _fmt_dollar(f.free_cash_flow)))
        if f.cash_conversion_ratio: metrics.append(("Cash Conv.", f"{f.cash_conversion_ratio:.0%}"))
        if f.net_debt_to_ltm_ebitda: metrics.append(("ND/EBITDA", _fmt_ratio(f.net_debt_to_ltm_ebitda)))

    if metrics:
        _add_metric_row(slide, metrics[:6], Inches(1.5))

    # WC trend table
    if analysis.working_capital and len(analysis.working_capital.periods) > 1:
        wc_periods = analysis.working_capital.periods[-6:]
        header = ["Period"] + [p.period for p in wc_periods]
        rows = [header]
        for label, attr in [("DSO", "dso"), ("DPO", "dpo"), ("CCC", "ccc")]:
            row = [label]
            for p in wc_periods:
                v = getattr(p, attr, None)
                row.append(f"{v:.0f}" if v else "—")
            rows.append(row)
        _add_table(slide, rows, Inches(0.5), Inches(3.5), Inches(10), [1.2] + [1.3] * len(wc_periods), font_size=10)


def _create_vc_plan(prs, vc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Value Creation Plan",
               f"Total Annual EBITDA Opportunity: {_fmt_dollar(vc.total_ebitda_opportunity)}")

    header = ["#", "Initiative", "Category", "Annual Impact", "Cost", "Timeline", "Confidence"]
    rows = [header]
    for i, init in enumerate(vc.prioritized_plan[:8], 1):
        rows.append([
            str(i),
            init.name[:40],
            init.category[:15],
            _fmt_dollar(init.ebitda_impact_annual),
            _fmt_dollar(init.implementation_cost),
            f"{init.timeline_months}mo",
            init.confidence,
        ])

    _add_table(slide, rows, Inches(0.3), Inches(1.8), Inches(12.5), [0.4, 4.0, 1.5, 1.5, 1.2, 0.9, 1.0], font_size=10)


def _create_ai_roadmap(prs, vc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "AI Transformation Roadmap")

    y = Inches(1.5)

    # Automation
    if vc.ai_automation_opportunities:
        _add_text_block(slide, "AI Automation Opportunities", Inches(0.5), y, Inches(12), Inches(0.3), font_size=16, color=DARK_BLUE)
        y = Inches(y.inches + 0.4)
        items = [f"{a.name} — {_fmt_dollar(a.ebitda_impact_annual)}/yr" + (f" ({', '.join(a.specific_tools)})" if a.specific_tools else "") for a in vc.ai_automation_opportunities[:4]]
        _add_bullet_list(slide, items, Inches(0.5), y, Inches(12), Inches(1.5), font_size=11)
        y = Inches(y.inches + 1.6)

    # Product
    if vc.ai_product_recommendations:
        _add_text_block(slide, "Product AI Features to Build", Inches(0.5), y, Inches(12), Inches(0.3), font_size=16, color=DARK_BLUE)
        y = Inches(y.inches + 0.4)
        _add_bullet_list(slide, vc.ai_product_recommendations[:3], Inches(0.5), y, Inches(12), Inches(1.2), font_size=11)
        y = Inches(y.inches + 1.3)

    # Risks
    if vc.ai_disruption_risks:
        _add_text_block(slide, "AI Disruption Risks", Inches(0.5), y, Inches(12), Inches(0.3), font_size=16, color=RED)
        y = Inches(y.inches + 0.4)
        _add_bullet_list(slide, vc.ai_disruption_risks[:3], Inches(0.5), y, Inches(12), Inches(1.2), font_size=11)


def _create_strategy_risks(prs, vc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Strategic Priorities & Risks")

    # Priorities on left
    if vc.strategic_priorities:
        _add_text_block(slide, "Strategic Priorities", Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.3), font_size=16, color=DARK_BLUE)
        _add_bullet_list(slide, vc.strategic_priorities[:5], Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.5), font_size=11)

    # Risks on right
    if vc.key_risks:
        _add_text_block(slide, "Key Risks", Inches(7), Inches(1.3), Inches(5.5), Inches(0.3), font_size=16, color=RED)
        _add_bullet_list(slide, vc.key_risks[:5], Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), font_size=11)


def _create_returns(prs, returns):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Returns & Exit Analysis")

    metrics = []
    if returns.moic: metrics.append(("MOIC", _fmt_ratio(returns.moic)))
    if returns.irr: metrics.append(("IRR", f"{returns.irr:.0%}"))
    metrics.append(("Holding Period", f"{returns.holding_period_years:.0f} years"))
    if returns.exit_ev: metrics.append(("Exit EV", _fmt_dollar(returns.exit_ev)))
    if returns.exit_equity: metrics.append(("Exit Equity", _fmt_dollar(returns.exit_equity)))

    _add_metric_row(slide, metrics[:5], Inches(1.5))

    # Returns table
    header = ["Metric", "Value"]
    rows = [header]
    rows.append(["Entry Equity", _fmt_dollar(returns.entry_equity)])
    if returns.exit_ebitda: rows.append(["Exit EBITDA", _fmt_dollar(returns.exit_ebitda)])
    rows.append(["Exit Multiple", _fmt_ratio(returns.exit_multiple)])
    if returns.exit_ev: rows.append(["Exit Enterprise Value", _fmt_dollar(returns.exit_ev)])
    if returns.net_debt_at_exit is not None: rows.append(["Net Debt at Exit", _fmt_dollar(returns.net_debt_at_exit)])
    if returns.exit_equity: rows.append(["Exit Equity Value", _fmt_dollar(returns.exit_equity)])

    _add_table(slide, rows, Inches(0.5), Inches(3.5), Inches(5), [3.0, 2.0])
